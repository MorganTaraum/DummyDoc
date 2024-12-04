import os
import zipfile
from docx import Document
from pptx import Presentation
from openpyxl import Workbook
from fpdf import FPDF
import random
import requests
import time

Header = ["Account Opening Form", 
"Customer Transaction History", 
"Monthly Account Statement", 
"Loan Application Form", 
"Fixed Deposit Receipt", 
"Cheque Clearance Report", 
"Interest Rate Analysis", 
"Customer Complaint Log", 
"Branch Performance Overview", 
"Cash Withdrawal Limit Policy", 
"Financial Analysis & Reporting", 
"Profit and Loss Statement", 
"Quarterly Financial Report", 
"Asset and Liability Management", 
"Risk Assessment Report", 
"Liquidity Coverage Ratio Analysis", 
"Credit Risk Analysis", 
"Revenue Projections", 
"Operational Expense Report", 
"Capital Adequacy Report", 
"Annual Budget Proposal", 
"Compliance and Legal", 
"Anti-Money Laundering Policy", 
"KYC (Know Your Customer) Checklist", 
"Regulatory Compliance Report", 
"Internal Audit Findings", 
"Fraud Investigation Summary", 
"Data Privacy Policy", 
"Customer Confidentiality Agreement", 
"Central Bank Reporting Guidelines", 
"Sanction Screening Procedure", 
"Legal Notices and Correspondence", 
"Customer Service", 
"Customer Feedback Analysis", 
"Call Center Performance Metrics", 
"Service Request Log", 
"Loan Repayment Schedule", 
"Credit Card Usage Summary", 
"Customer Retention Strategy", 
"Priority Banking Customer List", 
"VIP Customer Interaction Log", 
"Query Resolution Time Analysis", 
"Online Banking User Statistics", 
"Training and Presentations", 
"Staff Training Manual", 
"Cybersecurity Awareness Guide", 
"New Product Launch Strategy", 
"Team Performance Review", 
"Marketing Campaign Presentation", 
"Banking Software Tutorial", 
"Employee Performance Goals", 
"Leadership Development Plan", 
"Banking Industry Trends Overview", 
"Risk Management Workshop Slides", ]

ITHeader = ["Operating System Configuration Backup", 
"Network Interface Settings", 
"Firewall Rules Configuration", 
"Database Server Configuration", 
"Proxy Server Settings", 
"Load Balancer Configuration Details", 
"VPN Access Configuration", 
"Active Directory Backup File", 
"Security Policy Enforcement Script", 
"Kernel Module Settings", 
"Logs and Monitoring", 
"System Event Log", 
"Authentication Log", 
"Failed Login Attempts Report", 
"Audit Log Archive", 
"Debugging Log File", 
"Performance Monitoring Metrics", 
"Network Traffic Analysis Report", 
"Error Handling Log", 
"Application Crash Dump", 
"Real-Time Monitoring Alerts", 
"Binary Files and Executables", 
"Core System Binary", 
"System Restore Point File", 
"Disaster Recovery Plan Binary", 
"Incremental Backup Archive", 
"Full System Image Backup", 
"Database Dump File", 
"Configuration Snapshot File", 
"Archived Log Binary File", 
"Secure Password Vault Backup", 
"Cloud Data Backup Log", 
"Redundant Array Backup Script", 
"Security and Compliance", 
"Encrypted Key Store File", 
"Access Control Policy Binary", 
"Patch Management Script", 
"Anti-Malware Signature File", 
"Security Patch Binary", 
"System Vulnerability Assessment Report", 
"File Integrity Monitoring Report", 
"Privileged Access Session Log"]

def RandomHeader():
    res = random.choice(Header)
    Header.remove(res)
    return res

def RandomIT():
    res = random.choice(ITHeader)
    ITHeader.remove(res)
    return res

def RandomText():
    time.sleep(3)
    paragraph = random.randint(0,10)
    len = random.choice(["short", "medium"])
    link = random.choice(["/link", ""])
    url = "https://loripsum.net/api/" + str(paragraph) + link + f"/{len}/plaintext"
    res = requests.get(url)
    #print(res.content)
    return (res.content).decode()

def create_dummy_files(output_folder, num_files):
    os.makedirs(output_folder, exist_ok=True)

    for i in range(1, num_files + 1):
        # Generate .docx file
        HeaderDoc = RandomHeader()
        doc = Document()
        doc.add_heading(f"{HeaderDoc} 2024", level=1)
        doc.add_paragraph(RandomText())
        doc.save(os.path.join(output_folder, f"{HeaderDoc}.docx"))
        print("Created " + HeaderDoc + ".docx")

        # Generate .xlsx file
        xlsxHeader = RandomHeader()
        wb = Workbook()
        ws = wb.active
        ws.title = f"Sheet{i}"
        ws.append(["ID", "Name", "Value"])
        ws.append([1, "Example", random.randint(0, 100000)])
        wb.save(os.path.join(output_folder, f"{xlsxHeader}.xlsx"))
        print("Created " + xlsxHeader + ".xlsx")

        # Generate .zip file

        zipHeader = RandomHeader()
        zip_path = os.path.join(output_folder, f"{zipHeader}.zip")
        with zipfile.ZipFile(zip_path, "w") as zf:
            zf.writestr(f"{zipHeader}.txt", RandomText())
        print("Created " + zipHeader + ".zip")

        # Generate .pptx file
        pptHeader = RandomHeader()
        ppt = Presentation()
        slide = ppt.slides.add_slide(ppt.slide_layouts[0])
        slide.shapes.title.text = f"{pptHeader}"
        slide.placeholders[1].text = RandomText()
        ppt.save(os.path.join(output_folder, f"{pptHeader}.pptx"))
        
        print("Created " + pptHeader + ".pptx")

        # Generate .pdf file
        pdf = FPDF()
        pdfHeader = RandomHeader()
        pdf.add_page()
        pdf.set_font("Arial", size=12)
        pdf.cell(200, 10, txt=f"{pdfHeader} 2024", ln=True)
        pdf.cell(200, 10, txt=RandomText(), ln=True)
        pdf.output(os.path.join(output_folder, f"{pdfHeader}.pdf"))

        print("Created " + pdfHeader + ".pdf")

        # Generate .txt file
        txtHeader = RandomIT()
        with open(os.path.join(output_folder, f"{txtHeader}.txt"), "w") as txt_file:
            txt_file.write(RandomText())

        print("Created " + txtHeader + ".txt")

        # Generate .bak file

        bakHeader = RandomIT()
        with open(os.path.join(output_folder, f"{bakHeader}.bak"), "w") as bak_file:
            bak_file.write(RandomText())

        print("Created " + bakHeader + ".bak")

    print(f"Created {num_files} files of each type in {output_folder}")


# Parameters
output_folder = "dummy_files"
number_of_files = 10  # Change this to create more or fewer files

# Generate files
create_dummy_files(output_folder, number_of_files)
